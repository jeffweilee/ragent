"""TDD — _PptxASTSplitter: PPTX binary → Document atoms (one per slide)."""

from __future__ import annotations

import io


def _make_pptx_bytes(slides: list[list[str]]) -> bytes:
    """Build a minimal PPTX. slides is a list of text-block lists per slide."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout

    for texts in slides:
        slide = prs.slides.add_slide(blank_layout)
        top = Inches(1)
        for text in texts:
            txBox = slide.shapes.add_textbox(Inches(1), top, Inches(6), Inches(0.5))
            txBox.text_frame.text = text
            top += Inches(0.6)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _run_splitter(data: bytes) -> list:
    from haystack.dataclasses import Document as HDoc

    from ragent.pipelines.factory import _PptxASTSplitter

    splitter = _PptxASTSplitter()
    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    doc = HDoc(content=None, meta={"mime_type": mime, "document_id": "doc-2", "raw_bytes": data})
    return splitter.run([doc])["documents"]


# ---------------------------------------------------------------------------
# Slide-level atoms
# ---------------------------------------------------------------------------


def test_pptx_one_atom_per_slide():
    data = _make_pptx_bytes([["Slide 1 title", "Slide 1 body"], ["Slide 2 title"]])
    atoms = _run_splitter(data)
    assert len(atoms) == 2


def test_pptx_atom_content_contains_slide_text():
    data = _make_pptx_bytes([["Hello World", "More text"]])
    atoms = _run_splitter(data)
    assert len(atoms) == 1
    assert "Hello World" in atoms[0].content
    assert "More text" in atoms[0].content


def test_pptx_raw_content_set():
    data = _make_pptx_bytes([["Slide text"]])
    atoms = _run_splitter(data)
    assert "raw_content" in atoms[0].meta
    assert "Slide text" in atoms[0].meta["raw_content"]


def test_pptx_meta_passthrough():
    data = _make_pptx_bytes([["text"]])
    atoms = _run_splitter(data)
    assert atoms[0].meta["document_id"] == "doc-2"
    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert atoms[0].meta["mime_type"] == mime


def test_pptx_empty_presentation_yields_no_atoms():
    from pptx import Presentation

    prs = Presentation()
    buf = io.BytesIO()
    prs.save(buf)
    atoms = _run_splitter(buf.getvalue())
    assert atoms == []


def test_pptx_blank_slide_skipped():
    """A slide with no text shapes produces no atom."""
    from pptx import Presentation

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    prs.slides.add_slide(blank_layout)  # no text added
    buf = io.BytesIO()
    prs.save(buf)
    atoms = _run_splitter(buf.getvalue())
    assert atoms == []


def test_pptx_slide_number_in_meta():
    """Each atom should carry slide_number (1-based) in meta."""
    data = _make_pptx_bytes([["First"], ["Second"], ["Third"]])
    atoms = _run_splitter(data)
    assert len(atoms) == 3
    slide_nums = [a.meta.get("slide_number") for a in atoms]
    assert slide_nums == [1, 2, 3]


def test_pptx_multiple_text_boxes_merged():
    """All text on a slide merges into a single atom."""
    data = _make_pptx_bytes([["Alpha", "Beta", "Gamma"]])
    atoms = _run_splitter(data)
    assert len(atoms) == 1
    for word in ["Alpha", "Beta", "Gamma"]:
        assert word in atoms[0].content
